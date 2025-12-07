#!/usr/bin/env python3
"""
PPT 디자인 통일화 스크립트
- 원본 슬라이드 1-9의 스타일을 분석하여 슬라이드 10-20에 적용
- 논리적 순서(순서/방법/시간)에 맞게 재구성
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from copy import deepcopy
import os

# 스타일 상수 정의 (원본 PPT 1-9 기반)
STYLE = {
    # 색상
    'TITLE_COLOR': RGBColor(0x1E, 0x40, 0xAF),  # #1E40AF - 메인 타이틀
    'SUBTITLE_COLOR': RGBColor(0x33, 0x41, 0x55),  # #334155 - 서브타이틀
    'BODY_COLOR': RGBColor(0x47, 0x55, 0x69),  # #475569 - 본문
    'ACCENT_COLOR': RGBColor(0x1E, 0x29, 0x3B),  # #1E293B - 강조
    'LIGHT_TEXT': RGBColor(0x64, 0x74, 0x8B),  # #64748B - 연한 텍스트
    'WHITE': RGBColor(0xFF, 0xFF, 0xFF),  # 흰색
    'GREEN': RGBColor(0x16, 0xA3, 0x4A),  # #16A34A - 녹색 강조
    
    # 폰트 크기
    'MAIN_TITLE_SIZE': Pt(27),
    'SUBTITLE_SIZE': Pt(13),
    'SECTION_TITLE_SIZE': Pt(15),
    'BODY_SIZE': Pt(12),
    'SMALL_SIZE': Pt(11),
    'TINY_SIZE': Pt(10),
    
    # 레이아웃 위치
    'TITLE_LEFT': Inches(0.5),
    'TITLE_TOP': Inches(0.5),
    'TITLE_WIDTH': Inches(12.33),
    'TITLE_HEIGHT': Inches(0.6),
    
    'SUBTITLE_TOP': Inches(1.1),
    'SUBTITLE_HEIGHT': Inches(0.4),
    
    'CONTENT_LEFT': Inches(0.5),
    'CONTENT_TOP': Inches(1.7),
    'CONTENT_WIDTH': Inches(12.33),
}

def create_title_shape(slide, title_text, subtitle_text=None):
    """원본 스타일의 제목 생성"""
    # 메인 타이틀
    title_shape = slide.shapes.add_textbox(
        STYLE['TITLE_LEFT'], STYLE['TITLE_TOP'],
        STYLE['TITLE_WIDTH'], STYLE['TITLE_HEIGHT']
    )
    title_frame = title_shape.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title_text
    p.font.size = STYLE['MAIN_TITLE_SIZE']
    p.font.bold = True
    p.font.color.rgb = STYLE['TITLE_COLOR']
    
    # 서브타이틀
    if subtitle_text:
        subtitle_shape = slide.shapes.add_textbox(
            STYLE['TITLE_LEFT'], STYLE['SUBTITLE_TOP'],
            STYLE['TITLE_WIDTH'], STYLE['SUBTITLE_HEIGHT']
        )
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle_text
        p.font.size = STYLE['SUBTITLE_SIZE']
        p.font.color.rgb = STYLE['SUBTITLE_COLOR']
    
    return title_shape

def create_insight_box(slide, title, content, left, top, width=Inches(3.5), height=Inches(1.0)):
    """인사이트 박스 생성 (슬라이드 5, 8 스타일)"""
    # 제목
    title_shape = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = STYLE['ACCENT_COLOR']
    
    # 내용
    content_shape = slide.shapes.add_textbox(left, top + Inches(0.3), width, height - Inches(0.3))
    tf = content_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = STYLE['BODY_SIZE']
    p.font.color.rgb = STYLE['BODY_COLOR']
    
    return title_shape, content_shape

def create_key_finding_box(slide, text, left, top, width=Inches(12)):
    """핵심 발견 박스 생성 (슬라이드 5, 8 하단 스타일)"""
    shape = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = STYLE['TITLE_COLOR']
    return shape

def create_bullet_point(slide, text, left, top, width=Inches(6)):
    """불릿 포인트 생성"""
    shape = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"▶ {text}"
    p.font.size = STYLE['SMALL_SIZE']
    p.font.color.rgb = STYLE['BODY_COLOR']
    return shape

def redesign_slide_13(slide):
    """슬라이드 13: 연도별 토지유형 변화량 분석 - 재디자인"""
    # 기존 shape들의 정보 수집 (이미지 제외)
    image_shapes = []
    text_shapes = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture
            image_shapes.append({
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            })
        elif shape.has_text_frame:
            text_shapes.append(shape)
    
    # 텍스트 shape 스타일 수정
    for shape in text_shapes:
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if '연도별 토지유형' in text or '변화량 분석' in text:
                # 메인 타이틀
                para.font.size = STYLE['MAIN_TITLE_SIZE']
                para.font.bold = True
                para.font.color.rgb = STYLE['TITLE_COLOR']
            elif '임야/농경지' in text or '추세 확인' in text:
                # 서브타이틀
                para.font.size = STYLE['SUBTITLE_SIZE']
                para.font.color.rgb = STYLE['SUBTITLE_COLOR']
            elif text.startswith('▶'):
                # 불릿 포인트
                para.font.size = STYLE['SMALL_SIZE']
                para.font.color.rgb = STYLE['BODY_COLOR']

def redesign_analysis_slide(slide, is_correlation=False):
    """분석 슬라이드 재디자인 (슬라이드 14-20)"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                # 타이틀 식별 및 스타일 적용
                title_keywords = ['PCA', 'K-Means', '군집화', '청주시', '인구', '도로', 
                                '공장용지', '상관관계', '변화', '분석', '비교', '성장']
                is_title = any(kw in text for kw in title_keywords) and len(text) < 50
                
                if is_title and ('분석' in text or '검증' in text or '비교' in text or 
                               '상관관계' in text or '성장' in text or '추이' in text):
                    # 메인 타이틀
                    para.font.size = STYLE['MAIN_TITLE_SIZE']
                    para.font.bold = True
                    para.font.color.rgb = STYLE['TITLE_COLOR']
                elif text.startswith('▶'):
                    # 불릿 포인트 (핵심 발견)
                    para.font.size = STYLE['SMALL_SIZE']
                    para.font.color.rgb = STYLE['BODY_COLOR']
                elif any(kw in text for kw in ['확인', '유의', '뚜렷', '높음', '낮음']) and len(text) < 80:
                    # 서브타이틀
                    para.font.size = STYLE['SUBTITLE_SIZE']
                    para.font.color.rgb = STYLE['SUBTITLE_COLOR']

def main():
    """메인 실행 함수"""
    # PPT 파일 경로
    input_pptx = "result_v3/산업 빅데이터 분석 실제 프로젝트결과서-충청북도 지적통계 토지 이용 현황 분석(최종)_v2.pptx"
    output_pptx = "result_v3/산업 빅데이터 분석 실제 프로젝트결과서-충청북도 지적통계 토지 이용 현황 분석(최종)_v3.pptx"
    
    print("PPT 로딩 중...")
    prs = Presentation(input_pptx)
    
    print(f"총 슬라이드 수: {len(prs.slides)}")
    
    # 슬라이드 목록 가져오기
    slides_list = list(prs.slides._sldIdLst)
    
    # 슬라이드 10-20 재디자인
    print("\n슬라이드 재디자인 중...")
    for idx, sldId in enumerate(slides_list[12:], 13):  # 13번부터 (인덱스 12)
        try:
            slide = prs.slides.get(sldId.id)
            if slide:
                print(f"  슬라이드 {idx} 재디자인 중...")
                redesign_analysis_slide(slide)
        except Exception as e:
            print(f"  슬라이드 {idx} 오류: {e}")
    
    # 저장
    print(f"\n저장 중: {output_pptx}")
    prs.save(output_pptx)
    print("완료!")
    
    return output_pptx

if __name__ == "__main__":
    main()
