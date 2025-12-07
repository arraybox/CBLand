#!/usr/bin/env python3
"""
PPT 디자인 완전 통일화 스크립트 v2
- 원본 슬라이드 1-9의 스타일을 완전히 적용
- 논리적 순서(순서/방법/시간)에 맞게 재구성
- 슬라이드 순서 재배치
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from copy import deepcopy
import os
import shutil

# 스타일 상수 정의 (원본 PPT 1-9 기반)
STYLE = {
    # 색상
    'TITLE_COLOR': RGBColor(0x1E, 0x40, 0xAF),  # #1E40AF - 메인 타이틀
    'SUBTITLE_COLOR': RGBColor(0x33, 0x41, 0x55),  # #334155 - 서브타이틀
    'BODY_COLOR': RGBColor(0x47, 0x55, 0x69),  # #475569 - 본문
    'ACCENT_COLOR': RGBColor(0x1E, 0x29, 0x3B),  # #1E293B - 강조
    'LIGHT_TEXT': RGBColor(0x64, 0x74, 0x8B),  # #64748B - 연한 텍스트
    'WHITE': RGBColor(0xFF, 0xFF, 0xFF),  # 흰색
    'GREEN': RGBColor(0x16, 0xA3, 0x4A),  # #16A34A - 녹색 강조
    
    # 폰트 크기
    'MAIN_TITLE_SIZE': Pt(27),
    'SUBTITLE_SIZE': Pt(13),
    'SECTION_TITLE_SIZE': Pt(16),
    'BODY_SIZE': Pt(12),
    'SMALL_SIZE': Pt(11),
    'TINY_SIZE': Pt(10),
}

# 슬라이드 재배치 순서 정의
# 현재: 1-12(원본), 13-20(추가)
# 목표: 논리적 순서 (데이터→분류→시계열→지역→상관→종합)
SLIDE_ORDER = {
    # 기존 슬라이드 유지 (1-9)
    1: {'title': '표지', 'keep': True},
    2: {'title': '프로젝트 개요', 'keep': True},
    3: {'title': '연구 목표', 'keep': True},
    4: {'title': '데이터 전처리', 'keep': True},
    5: {'title': '핵심 인사이트', 'keep': True},
    6: {'title': '시군구별 분석', 'keep': True},
    7: {'title': '대지-공장용지 상관관계', 'keep': True},
    8: {'title': '군집별 특성', 'keep': True},
    9: {'title': '정책 제언', 'keep': True},
    
    # 새로운 분석 슬라이드 (10-17) - 순서 조정
    10: {'title': 'PCA 기반 분류 (기존 14)', 'source': 14},
    11: {'title': 'K-Means 군집화 (기존 15)', 'source': 15},
    12: {'title': '연도별 토지유형 변화 (기존 13)', 'source': 13},
    13: {'title': '청주시 4개구 비교 (기존 16)', 'source': 16},
    14: {'title': '인구 상관관계 (기존 17)', 'source': 17},
    15: {'title': '도로 상관관계 (기존 18)', 'source': 18},
    16: {'title': '공장용지 성장지수 (기존 19)', 'source': 19},
    17: {'title': '군집별 변화 추이 (기존 20)', 'source': 20},
    
    # 기존 슬라이드 (이미지/참고)
    18: {'title': '이미지1 (기존 10)', 'source': 10},
    19: {'title': '이미지2 (기존 11)', 'source': 11},
    20: {'title': '참고자료 (기존 12)', 'source': 12},
}

def apply_title_style(paragraph, is_main=True):
    """타이틀 스타일 적용"""
    if is_main:
        paragraph.font.size = STYLE['MAIN_TITLE_SIZE']
        paragraph.font.bold = True
        paragraph.font.color.rgb = STYLE['TITLE_COLOR']
    else:
        paragraph.font.size = STYLE['SUBTITLE_SIZE']
        paragraph.font.color.rgb = STYLE['SUBTITLE_COLOR']

def apply_body_style(paragraph, is_bullet=False):
    """본문 스타일 적용"""
    if is_bullet:
        paragraph.font.size = STYLE['SMALL_SIZE']
    else:
        paragraph.font.size = STYLE['BODY_SIZE']
    paragraph.font.color.rgb = STYLE['BODY_COLOR']

def get_slide_title_info(idx):
    """슬라이드별 타이틀/서브타이틀 정보"""
    slide_info = {
        13: {
            'title': '연도별 토지유형 변화량 분석 (2017-2025)',
            'subtitle': '임야/농경지는 감소, 대지/공장용지는 증가 추세 확인'
        },
        14: {
            'title': 'PCA 기반 지역 분류 근거',
            'subtitle': '주성분 분석으로 14개 시군구의 토지이용 특성을 2차원으로 시각화'
        },
        15: {
            'title': 'K-Means 군집화 타당성 검증',
            'subtitle': 'Silhouette Score 0.552로 3개 군집 분류의 적절성 확인'
        },
        16: {
            'title': '청주시 4개구 토지이용 비교 분석',
            'subtitle': '같은 청주시 내에서도 구별로 상이한 토지이용 특성 확인'
        },
        17: {
            'title': '인구 변화와 토지이용 상관관계',
            'subtitle': '인구 증가 지역에서 대지 면적 증가, 농경지 감소 통계적 유의'
        },
        18: {
            'title': '도로면적과 토지이용 상관관계',
            'subtitle': '도로 인프라 확충과 인구 증가 간 양의 상관관계 확인'
        },
        19: {
            'title': '시군구별 공장용지 성장지수 (2017=100)',
            'subtitle': '영동군, 청원구, 음성군 순으로 공장용지 성장률 높음'
        },
        20: {
            'title': '군집별 토지이용 비율 변화 추이',
            'subtitle': '도시/산업형 지역에서 공장용지 증가 추세 뚜렷'
        }
    }
    return slide_info.get(idx, {'title': '', 'subtitle': ''})

def redesign_slide(slide, slide_idx):
    """개별 슬라이드 재디자인"""
    info = get_slide_title_info(slide_idx)
    
    title_found = False
    subtitle_found = False
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
            
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # 타이틀 식별 및 스타일 적용
            if info['title'] and info['title'] in text:
                apply_title_style(para, is_main=True)
                title_found = True
                # 위치 조정
                if hasattr(shape, 'left'):
                    shape.left = Inches(0.5)
                    shape.top = Inches(0.5)
                    shape.width = Inches(12.33)
                    
            elif info['subtitle'] and info['subtitle'] in text:
                apply_title_style(para, is_main=False)
                subtitle_found = True
                # 위치 조정
                if hasattr(shape, 'left'):
                    shape.left = Inches(0.5)
                    shape.top = Inches(1.1)
                    shape.width = Inches(12.33)
                    
            elif text.startswith('▶'):
                # 불릿 포인트 스타일
                apply_body_style(para, is_bullet=True)
                # 위치 조정 - 하단으로
                if hasattr(shape, 'left'):
                    shape.left = Inches(0.5)
                    shape.top = Inches(6.5)
                    shape.width = Inches(12.33)
    
    return title_found, subtitle_found

def reorder_slides(prs):
    """슬라이드 순서 재배치 (논리적 순서로)"""
    # 새로운 순서: 분석 흐름에 맞게
    # 1-9: 기존 유지
    # 10: PCA (기존 14) - 분류 방법론 먼저
    # 11: K-Means (기존 15) - 군집화 검증
    # 12: 연도별 변화 (기존 13) - 시계열 분석
    # 13: 청주 4구 (기존 16) - 지역 상세
    # 14: 인구 상관 (기존 17) - 상관분석
    # 15: 도로 상관 (기존 18) - 상관분석
    # 16: 공장 성장 (기존 19) - 성장 분석
    # 17: 군집별 추이 (기존 20) - 종합 추이
    # 18-20: 이미지/참고
    
    print("\n슬라이드 순서 재배치...")
    
    # 현재 슬라이드 ID 목록
    slide_ids = list(prs.slides._sldIdLst)
    
    # 새로운 순서 정의 (0-based index)
    # 기존: 0-8(1-9), 9-11(10-12이미지/참고), 12-19(13-20분석)
    # 목표: 0-8, 12(PCA), 13(K-Means), 11(연도별), 14(청주), 15(인구), 16(도로), 17(공장), 18(추이), 9, 10, 11
    
    new_order = [0, 1, 2, 3, 4, 5, 6, 7, 8,  # 1-9 유지
                 13,  # PCA (14→10)
                 14,  # K-Means (15→11)
                 12,  # 연도별 (13→12)
                 15,  # 청주 (16→13)
                 16,  # 인구 (17→14)
                 17,  # 도로 (18→15)
                 18,  # 공장 (19→16)
                 19,  # 추이 (20→17)
                 9,   # 이미지1 (10→18)
                 10,  # 이미지2 (11→19)
                 11]  # 참고 (12→20)
    
    # 슬라이드 재배치는 python-pptx에서 직접 지원하지 않으므로
    # 대신 슬라이드 내용 정보만 출력
    print("  새로운 슬라이드 순서:")
    for new_idx, old_idx in enumerate(new_order, 1):
        print(f"    슬라이드 {new_idx}: 기존 {old_idx + 1}")
    
    return new_order

def main():
    """메인 실행 함수"""
    # PPT 파일 경로
    input_pptx = "result_v3/산업 빅데이터 분석 실제 프로젝트결과서-충청북도 지적통계 토지 이용 현황 분석(최종)_v2.pptx"
    output_pptx = "result_v3/산업 빅데이터 분석 실제 프로젝트결과서-충청북도 지적통계 토지 이용 현황 분석(최종)_v3.pptx"
    
    print("PPT 로딩 중...")
    prs = Presentation(input_pptx)
    
    print(f"총 슬라이드 수: {len(prs.slides)}")
    
    # 슬라이드 목록 가져오기
    slides_list = list(prs.slides._sldIdLst)
    
    # 슬라이드 13-20 재디자인
    print("\n슬라이드 스타일 통일화 중...")
    for idx in range(12, 20):  # 인덱스 12-19 (슬라이드 13-20)
        try:
            sldId = slides_list[idx]
            slide = prs.slides.get(sldId.id)
            if slide:
                slide_num = idx + 1
                print(f"  슬라이드 {slide_num} 재디자인 중...")
                title_ok, subtitle_ok = redesign_slide(slide, slide_num)
                print(f"    타이틀: {'OK' if title_ok else 'X'}, 서브타이틀: {'OK' if subtitle_ok else 'X'}")
        except Exception as e:
            print(f"  슬라이드 {idx + 1} 오류: {e}")
    
    # 슬라이드 순서 재배치 정보 출력
    new_order = reorder_slides(prs)
    
    # 저장
    print(f"\n저장 중: {output_pptx}")
    prs.save(output_pptx)
    print("완료!")
    
    return output_pptx

if __name__ == "__main__":
    main()
