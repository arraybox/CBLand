#!/usr/bin/env python3
"""
PPT 디자인 완전 통일화 및 순서 재배치 최종 스크립트
- 원본 슬라이드 1-9의 스타일을 완전히 적용
- 논리적 순서(순서/방법/시간)에 맞게 슬라이드 실제 재배치
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy
from lxml import etree
import os
import shutil

# 스타일 상수 정의 (원본 PPT 1-9 기반)
STYLE = {
    'TITLE_COLOR': RGBColor(0x1E, 0x40, 0xAF),
    'SUBTITLE_COLOR': RGBColor(0x33, 0x41, 0x55),
    'BODY_COLOR': RGBColor(0x47, 0x55, 0x69),
    'ACCENT_COLOR': RGBColor(0x1E, 0x29, 0x3B),
    'LIGHT_TEXT': RGBColor(0x64, 0x74, 0x8B),
    'WHITE': RGBColor(0xFF, 0xFF, 0xFF),
    'GREEN': RGBColor(0x16, 0xA3, 0x4A),
    'MAIN_TITLE_SIZE': Pt(27),
    'SUBTITLE_SIZE': Pt(13),
    'SECTION_TITLE_SIZE': Pt(16),
    'BODY_SIZE': Pt(12),
    'SMALL_SIZE': Pt(11),
    'TINY_SIZE': Pt(10),
}

def move_slide(presentation, old_index, new_index):
    """슬라이드를 old_index에서 new_index로 이동"""
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    
    # 이동할 슬라이드 추출
    slide_to_move = slides[old_index]
    
    # 원래 위치에서 제거
    xml_slides.remove(slide_to_move)
    
    # 새 위치에 삽입
    if new_index >= len(xml_slides):
        xml_slides.append(slide_to_move)
    else:
        xml_slides.insert(new_index, slide_to_move)

def reorder_slides_actual(prs):
    """실제로 슬라이드 순서 재배치"""
    # 목표 순서:
    # 1-9: 유지
    # 10: PCA (기존 14) 
    # 11: K-Means (기존 15)
    # 12: 연도별 변화 (기존 13)
    # 13: 청주 4구 (기존 16)
    # 14: 인구 상관 (기존 17)
    # 15: 도로 상관 (기존 18)
    # 16: 공장 성장 (기존 19)
    # 17: 군집별 추이 (기존 20)
    # 18: 이미지1 (기존 10)
    # 19: 이미지2 (기존 11)
    # 20: 참고자료 (기존 12)
    
    print("\n슬라이드 순서 재배치 중...")
    
    # 순차적으로 이동 (복잡한 재배치를 위해 단계별로)
    # 현재: [0-8, 9(이미지1), 10(이미지2), 11(참고), 12(연도별), 13(PCA), 14(K-Means), 15(청주), 16(인구), 17(도로), 18(공장), 19(추이)]
    # 목표: [0-8, 13(PCA), 14(K-Means), 12(연도별), 15(청주), 16(인구), 17(도로), 18(공장), 19(추이), 9, 10, 11]
    
    # 단계 1: PCA(13)를 9번 위치로 이동
    move_slide(prs, 13, 9)
    print("  1단계: PCA 슬라이드 이동 완료")
    # 결과: [0-8, PCA, 9(이미지1), 10(이미지2), 11(참고), 12(연도별), 14(K-Means), 15(청주), 16(인구), 17(도로), 18(공장), 19(추이)]
    
    # 단계 2: K-Means(14, 현재 위치 14)를 10번 위치로 이동
    move_slide(prs, 14, 10)
    print("  2단계: K-Means 슬라이드 이동 완료")
    # 결과: [0-8, PCA, K-Means, 9(이미지1), 10(이미지2), 11(참고), 12(연도별), 15(청주), 16(인구), 17(도로), 18(공장), 19(추이)]
    
    # 단계 3: 연도별(14, 현재 위치 14)를 11번 위치로 이동
    move_slide(prs, 14, 11)
    print("  3단계: 연도별 변화 슬라이드 이동 완료")
    # 결과: [0-8, PCA, K-Means, 연도별, 9(이미지1), 10(이미지2), 11(참고), 15(청주), 16(인구), 17(도로), 18(공장), 19(추이)]
    
    # 단계 4: 청주(15, 현재 위치 15)를 12번 위치로 이동
    move_slide(prs, 15, 12)
    print("  4단계: 청주 4구 슬라이드 이동 완료")
    # 결과: [0-8, PCA, K-Means, 연도별, 청주, 9(이미지1), 10(이미지2), 11(참고), 16(인구), 17(도로), 18(공장), 19(추이)]
    
    # 단계 5: 인구(16, 현재 위치 16)를 13번 위치로 이동
    move_slide(prs, 16, 13)
    print("  5단계: 인구 상관 슬라이드 이동 완료")
    # 결과: [0-8, PCA, K-Means, 연도별, 청주, 인구, 9(이미지1), 10(이미지2), 11(참고), 17(도로), 18(공장), 19(추이)]
    
    # 단계 6: 도로(17, 현재 위치 17)를 14번 위치로 이동
    move_slide(prs, 17, 14)
    print("  6단계: 도로 상관 슬라이드 이동 완료")
    # 결과: [0-8, PCA, K-Means, 연도별, 청주, 인구, 도로, 9(이미지1), 10(이미지2), 11(참고), 18(공장), 19(추이)]
    
    # 단계 7: 공장(18, 현재 위치 18)를 15번 위치로 이동
    move_slide(prs, 18, 15)
    print("  7단계: 공장 성장 슬라이드 이동 완료")
    # 결과: [0-8, PCA, K-Means, 연도별, 청주, 인구, 도로, 공장, 9(이미지1), 10(이미지2), 11(참고), 19(추이)]
    
    # 단계 8: 추이(19, 현재 위치 19)를 16번 위치로 이동
    move_slide(prs, 19, 16)
    print("  8단계: 군집별 추이 슬라이드 이동 완료")
    # 최종: [0-8, PCA, K-Means, 연도별, 청주, 인구, 도로, 공장, 추이, 9(이미지1), 10(이미지2), 11(참고)]

def get_slide_title_info(slide_num):
    """슬라이드별 타이틀/서브타이틀 정보 (재배치 후 기준)"""
    slide_info = {
        10: {
            'title': 'PCA 기반 지역 분류 근거',
            'subtitle': '주성분 분석으로 14개 시군구의 토지이용 특성을 2차원으로 시각화'
        },
        11: {
            'title': 'K-Means 군집화 타당성 검증',
            'subtitle': 'Silhouette Score 0.552로 3개 군집 분류의 적절성 확인'
        },
        12: {
            'title': '연도별 토지유형 변화량 분석 (2017-2025)',
            'subtitle': '임야/농경지는 감소, 대지/공장용지는 증가 추세 확인'
        },
        13: {
            'title': '청주시 4개구 토지이용 비교 분석',
            'subtitle': '같은 청주시 내에서도 구별로 상이한 토지이용 특성 확인'
        },
        14: {
            'title': '인구 변화와 토지이용 상관관계',
            'subtitle': '인구 증가 지역에서 대지 면적 증가, 농경지 감소 통계적 유의'
        },
        15: {
            'title': '도로면적과 토지이용 상관관계',
            'subtitle': '도로 인프라 확충과 인구 증가 간 양의 상관관계 확인'
        },
        16: {
            'title': '시군구별 공장용지 성장지수 (2017=100)',
            'subtitle': '영동군, 청원구, 음성군 순으로 공장용지 성장률 높음'
        },
        17: {
            'title': '군집별 토지이용 비율 변화 추이',
            'subtitle': '도시/산업형 지역에서 공장용지 증가 추세 뚜렷'
        }
    }
    return slide_info.get(slide_num, {'title': '', 'subtitle': ''})

def redesign_slide(slide, slide_num):
    """개별 슬라이드 재디자인"""
    info = get_slide_title_info(slide_num)
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
            
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # 타이틀 매칭 (부분 매칭)
            title_keywords = ['PCA', 'K-Means', '군집화', '청주시', '인구', '도로', 
                            '공장용지', '상관관계', '변화', '분석', '비교', '성장',
                            '연도별', '토지유형', '추이', '검증', '근거', '타당성']
            
            is_title = False
            for kw in title_keywords:
                if kw in text and len(text) < 60:
                    is_title = True
                    break
            
            if is_title and len(text) < 60:
                # 메인 타이틀 스타일
                para.font.size = STYLE['MAIN_TITLE_SIZE']
                para.font.bold = True
                para.font.color.rgb = STYLE['TITLE_COLOR']
                # 위치 조정
                if hasattr(shape, 'left'):
                    shape.left = Inches(0.5)
                    shape.top = Inches(0.5)
                    shape.width = Inches(12.33)
                    
            elif any(kw in text for kw in ['확인', '유의', '뚜렷', '높음', '추세', 
                                           '시각화', '적절성', '특성 확인', '감소']) and len(text) < 100:
                # 서브타이틀 스타일
                para.font.size = STYLE['SUBTITLE_SIZE']
                para.font.color.rgb = STYLE['SUBTITLE_COLOR']
                if hasattr(shape, 'left'):
                    shape.left = Inches(0.5)
                    shape.top = Inches(1.1)
                    shape.width = Inches(12.33)
                    
            elif text.startswith('▶'):
                # 불릿 포인트 스타일
                para.font.size = STYLE['SMALL_SIZE']
                para.font.color.rgb = STYLE['BODY_COLOR']
                if hasattr(shape, 'left'):
                    shape.left = Inches(0.5)
                    shape.top = Inches(6.5)
                    shape.width = Inches(12.33)

def main():
    """메인 실행 함수"""
    # PPT 파일 경로
    input_pptx = "result_v3/산업 빅데이터 분석 실제 프로젝트결과서-충청북도 지적통계 토지 이용 현황 분석(최종)_v2.pptx"
    output_pptx = "result_v3/산업 빅데이터 분석 실제 프로젝트결과서-충청북도 지적통계 토지 이용 현황 분석(최종)_v3.pptx"
    
    print("=" * 60)
    print("PPT 디자인 통일화 및 순서 재배치")
    print("=" * 60)
    
    print("\n1. PPT 로딩 중...")
    prs = Presentation(input_pptx)
    print(f"   총 슬라이드 수: {len(prs.slides)}")
    
    # 슬라이드 순서 재배치
    print("\n2. 슬라이드 순서 재배치...")
    reorder_slides_actual(prs)
    
    # 슬라이드 목록 가져오기 (재배치 후)
    slides_list = list(prs.slides._sldIdLst)
    
    # 슬라이드 10-17 재디자인 (재배치 후 새로운 분석 슬라이드)
    print("\n3. 슬라이드 스타일 통일화 중...")
    for idx in range(9, 17):  # 인덱스 9-16 (슬라이드 10-17)
        try:
            sldId = slides_list[idx]
            slide = prs.slides.get(sldId.id)
            if slide:
                slide_num = idx + 1
                print(f"   슬라이드 {slide_num} 재디자인 중...")
                redesign_slide(slide, slide_num)
        except Exception as e:
            print(f"   슬라이드 {idx + 1} 오류: {e}")
    
    # 저장
    print(f"\n4. 저장 중: {output_pptx}")
    prs.save(output_pptx)
    
    print("\n" + "=" * 60)
    print("완료!")
    print("=" * 60)
    
    # 최종 슬라이드 순서 출력
    print("\n최종 슬라이드 순서:")
    final_order = [
        "1. 표지",
        "2. 프로젝트 개요",
        "3. 연구 목표",
        "4. 데이터 전처리",
        "5. 핵심 인사이트",
        "6. 시군구별 분석",
        "7. 대지-공장용지 상관관계",
        "8. 군집별 특성",
        "9. 정책 제언",
        "10. PCA 기반 지역 분류 근거 ★",
        "11. K-Means 군집화 타당성 검증 ★",
        "12. 연도별 토지유형 변화량 분석 ★",
        "13. 청주시 4개구 비교 분석 ★",
        "14. 인구-토지이용 상관관계 ★",
        "15. 도로-토지이용 상관관계 ★",
        "16. 공장용지 성장지수 ★",
        "17. 군집별 변화 추이 ★",
        "18. 이미지 1",
        "19. 이미지 2",
        "20. 참고자료"
    ]
    for item in final_order:
        print(f"  {item}")
    
    print("\n★ = 새로 추가된 분석 슬라이드 (스타일 통일화 완료)")
    
    return output_pptx

if __name__ == "__main__":
    main()
