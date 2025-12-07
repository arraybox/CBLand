#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT 업데이트 스크립트 v2
- 기존 1-9번 슬라이드 스타일과 통일
- 분석 순서/방법/시간에 맞게 적절한 위치에 삽입
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from copy import deepcopy
import os

# 경로 설정
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
RESULT_DIR = os.path.join(BASE_DIR, 'result_v3')
ORIGINAL_PPTX = os.path.join(BASE_DIR, 'result', '산업 빅데이터 분석 실제 프로젝트결과서-충청북도 지적통계 토지 이용 현황 분석(최종).pptx')
OUTPUT_PPTX = os.path.join(RESULT_DIR, '산업 빅데이터 분석 실제 프로젝트결과서-충청북도 지적통계 토지 이용 현황 분석(최종)_v2.pptx')

# 스타일 상수 (기존 PPT 분석 기반)
TITLE_FONT_SIZE = Pt(28)
SUBTITLE_FONT_SIZE = Pt(14)
BODY_FONT_SIZE = Pt(11)
INSIGHT_TITLE_SIZE = Pt(16)
INSIGHT_BODY_SIZE = Pt(11)

# 색상
COLOR_TITLE = RGBColor(0, 0, 0)
COLOR_SUBTITLE = RGBColor(100, 100, 100)
COLOR_BODY = RGBColor(50, 50, 50)
COLOR_HIGHLIGHT = RGBColor(200, 50, 50)
COLOR_ACCENT = RGBColor(0, 112, 192)

# 배경색
BG_LIGHT_GRAY = RGBColor(245, 245, 245)


def add_title_box(slide, text, left=0.4, top=0.3, width=10, height=0.5):
    """기존 스타일 제목 추가"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = TITLE_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    return shape


def add_subtitle_box(slide, text, left=0.4, top=0.8, width=10, height=0.4):
    """기존 스타일 부제목 추가"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = SUBTITLE_FONT_SIZE
    p.font.color.rgb = COLOR_SUBTITLE
    return shape


def add_separator_line(slide, left=0.4, top=1.2, width=0.6):
    """구분선 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()
    return shape


def add_insight_section(slide, title, body, left, top, icon_text="▶"):
    """인사이트 섹션 추가 (기존 스타일)"""
    # 아이콘/불릿
    icon_shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(0.3), Inches(0.3))
    icon_tf = icon_shape.text_frame
    icon_p = icon_tf.paragraphs[0]
    icon_p.text = icon_text
    icon_p.font.size = Pt(12)
    icon_p.font.color.rgb = COLOR_ACCENT
    
    # 제목
    title_shape = slide.shapes.add_textbox(Inches(left + 0.4), Inches(top), Inches(3), Inches(0.3))
    title_tf = title_shape.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = INSIGHT_TITLE_SIZE
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_TITLE
    
    # 본문
    body_shape = slide.shapes.add_textbox(Inches(left + 0.4), Inches(top + 0.35), Inches(5.5), Inches(0.5))
    body_tf = body_shape.text_frame
    body_tf.word_wrap = True
    body_p = body_tf.paragraphs[0]
    body_p.text = body
    body_p.font.size = INSIGHT_BODY_SIZE
    body_p.font.color.rgb = COLOR_BODY


def add_image_with_caption(slide, img_path, left, top, width, height, caption=None):
    """이미지와 캡션 추가"""
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))
        if caption:
            cap_shape = slide.shapes.add_textbox(
                Inches(left), Inches(top + height + 0.1), Inches(width), Inches(0.3)
            )
            cap_tf = cap_shape.text_frame
            cap_p = cap_tf.paragraphs[0]
            cap_p.text = caption
            cap_p.font.size = Pt(9)
            cap_p.font.color.rgb = COLOR_SUBTITLE
            cap_p.alignment = PP_ALIGN.CENTER
        return True
    return False


def add_header_banner(slide, text, top=0.0, height=1.2):
    """상단 배너 (기존 슬라이드 5, 8 스타일)"""
    # 배경 박스
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(top), Inches(13.33), Inches(height)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(0, 82, 147)  # 진한 파란색
    bg_shape.line.fill.background()
    
    # 텍스트
    txt_shape = slide.shapes.add_textbox(Inches(0.5), Inches(top + 0.3), Inches(12), Inches(0.5))
    txt_tf = txt_shape.text_frame
    txt_p = txt_tf.paragraphs[0]
    txt_p.text = text
    txt_p.font.size = Pt(22)
    txt_p.font.bold = True
    txt_p.font.color.rgb = RGBColor(255, 255, 255)


def add_footer_note(slide, text, top=7.1):
    """하단 주석"""
    shape = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12), Inches(0.3))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_SUBTITLE
    p.alignment = PP_ALIGN.CENTER


def create_slide_yearly_change(prs):
    """슬라이드: 연도별 토지유형 변화량 분석"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    
    # 제목
    add_title_box(slide, "연도별 토지유형 변화량 분석 (2017-2025)")
    add_separator_line(slide)
    add_subtitle_box(slide, "임야/농경지는 감소, 대지/공장용지는 증가 추세 확인", top=1.4)
    
    # 이미지
    img_path = os.path.join(RESULT_DIR, '05_yearly_landtype_change.png')
    add_image_with_caption(slide, img_path, 0.3, 1.9, 12.5, 4.8)
    
    # 인사이트
    add_insight_section(slide, "2017→2025 변화율", 
                       "임야 -0.84%, 농경지 -3.19%, 대지 +12.67%, 공장용지 +22.74%",
                       0.4, 6.85)
    
    return slide


def create_slide_pca_basis(prs):
    """슬라이드: PCA 분류 근거"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    
    add_title_box(slide, "PCA 기반 지역 분류 근거")
    add_separator_line(slide)
    add_subtitle_box(slide, "주성분 분석으로 14개 시군구의 토지이용 특성을 2차원으로 시각화", top=1.4)
    
    # 이미지
    img_path = os.path.join(RESULT_DIR, '02_PCA_classification_basis.png')
    add_image_with_caption(slide, img_path, 0.3, 1.9, 12.5, 4.5)
    
    # 인사이트
    insights = [
        ("PC1 (90.1%)", "도시/산업화 정도 - 공장용지/대지 비율 높을수록 양(+)의 방향"),
        ("PC2 (6.6%)", "농경지 비율 - 농경지가 높을수록 양(+)의 방향"),
        ("누적 설명력", "96.7% → 2개 주성분으로 토지이용 특성 대부분 설명")
    ]
    
    for i, (title, body) in enumerate(insights):
        add_insight_section(slide, title, body, 0.4, 6.5 + i * 0.35)
    
    return slide


def create_slide_silhouette(prs):
    """슬라이드: K-Means 실루엣 분석"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    
    add_title_box(slide, "K-Means 군집화 타당성 검증")
    add_separator_line(slide)
    add_subtitle_box(slide, "Silhouette Score 0.552로 3개 군집 분류의 적절성 확인", top=1.4)
    
    img_path = os.path.join(RESULT_DIR, '03_silhouette_analysis.png')
    add_image_with_caption(slide, img_path, 0.3, 1.9, 12.5, 4.3)
    
    insights = [
        ("Elbow Method", "k=3에서 최적 지점 확인"),
        ("Silhouette Score", "k=3에서 0.552 (0.4 이상이면 군집 분리 적절)"),
        ("군집별 분포", "모든 군집이 양(+)의 값으로 적절한 분류")
    ]
    
    for i, (title, body) in enumerate(insights):
        add_insight_section(slide, title, body, 0.4, 6.3 + i * 0.35)
    
    return slide


def create_slide_cheongju(prs):
    """슬라이드: 청주 4개구 비교"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    
    add_title_box(slide, "청주시 4개구 토지이용 비교 분석")
    add_separator_line(slide)
    add_subtitle_box(slide, "같은 청주시 내에서도 구별로 상이한 토지이용 특성 확인", top=1.4)
    
    img_path = os.path.join(RESULT_DIR, '09_cheongju_4gu_comparison.png')
    add_image_with_caption(slide, img_path, 0.3, 1.9, 12.5, 4.6)
    
    insights = [
        ("흥덕구/청원구", "도시/산업형 - 공장용지 비율 높음, 산업화 진행"),
        ("서원구", "균형형 - 도시/농촌 기능 혼재"),
        ("상당구", "농업/산림형 - 임야 76.7%로 산림 중심")
    ]
    
    for i, (title, body) in enumerate(insights):
        add_insight_section(slide, title, body, 0.4, 6.6 + i * 0.3)
    
    return slide


def create_slide_pop_correlation(prs):
    """슬라이드: 인구 상관관계"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    
    add_title_box(slide, "인구 변화와 토지이용 상관관계")
    add_separator_line(slide)
    add_subtitle_box(slide, "인구 증가 지역에서 대지 면적 증가, 농경지 감소가 통계적으로 유의미", top=1.4)
    
    img_path = os.path.join(RESULT_DIR, '10_correlation_population.png')
    add_image_with_caption(slide, img_path, 0.3, 1.9, 12.5, 4.5)
    
    insights = [
        ("인구↔대지", "r=0.602, p=0.023 (유의) → 인구 증가 지역에서 대지 면적 증가"),
        ("인구↔농경지", "r=-0.668, p=0.009 (유의) → 인구 증가 지역에서 농경지 감소"),
        ("인구↔공장용지", "r=-0.122, p=0.678 (비유의) → 공장용지는 인구와 직접적 연관 없음")
    ]
    
    for i, (title, body) in enumerate(insights):
        add_insight_section(slide, title, body, 0.4, 6.5 + i * 0.3)
    
    return slide


def create_slide_road_correlation(prs):
    """슬라이드: 도로 상관관계"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    
    add_title_box(slide, "도로면적과 토지이용 상관관계")
    add_separator_line(slide)
    add_subtitle_box(slide, "도로 인프라 확충과 인구 증가 간 양의 상관관계 확인", top=1.4)
    
    img_path = os.path.join(RESULT_DIR, '11_correlation_road.png')
    add_image_with_caption(slide, img_path, 0.3, 1.9, 12.5, 4.5)
    
    insights = [
        ("도로↔인구", "r=0.599, p=0.024 (유의) → 도로 인프라 확충 지역에서 인구 증가"),
        ("도로↔대지", "r=0.528, p=0.052 (경계) → 도로 확충과 대지 증가 연관성"),
        ("도로↔공장용지", "r=0.476, p=0.086 (비유의) → 공장용지는 산업단지 입지 특성 반영")
    ]
    
    for i, (title, body) in enumerate(insights):
        add_insight_section(slide, title, body, 0.4, 6.5 + i * 0.3)
    
    return slide


def create_slide_factory_heatmap(prs):
    """슬라이드: 공장용지 성장 히트맵"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    
    add_title_box(slide, "시군구별 공장용지 성장지수 (2017=100)")
    add_separator_line(slide)
    add_subtitle_box(slide, "영동군, 청원구, 음성군 순으로 공장용지 성장률 높음", top=1.4)
    
    img_path = os.path.join(RESULT_DIR, '12_heatmap_factory_growth.png')
    add_image_with_caption(slide, img_path, 0.5, 1.9, 12.0, 4.8)
    
    insights = [
        ("영동군", "+47.95% (1위) - 신규 산업단지 조성 영향"),
        ("청원구/음성군", "+34.05% / +25.62% - 기존 산업벨트 확장"),
        ("농업/산림형 지역", "보은, 단양 등도 공장용지 소폭 증가")
    ]
    
    for i, (title, body) in enumerate(insights):
        add_insight_section(slide, title, body, 0.4, 6.8 + i * 0.25)
    
    return slide


def create_slide_cluster_trend(prs):
    """슬라이드: 유형별 토지 변화 추이"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()
    
    add_title_box(slide, "군집별 토지이용 비율 변화 추이")
    add_separator_line(slide)
    add_subtitle_box(slide, "도시/산업형 지역에서 공장용지 증가 추세 뚜렷", top=1.4)
    
    img_path = os.path.join(RESULT_DIR, '06_landuse_trend_by_cluster.png')
    add_image_with_caption(slide, img_path, 0.3, 1.9, 12.5, 4.6)
    
    insights = [
        ("도시/산업형", "공장용지 지속 증가, 임야/농경지 감소 뚜렷"),
        ("농업/산림형", "변화 폭 미미, 기존 토지이용 구조 유지"),
        ("균형형", "대지 증가 추세, 도시화 진행 중")
    ]
    
    for i, (title, body) in enumerate(insights):
        add_insight_section(slide, title, body, 0.4, 6.6 + i * 0.3)
    
    return slide


def main():
    print("=" * 60)
    print("PPT 업데이트 v2 시작 (기존 스타일 통일)")
    print("=" * 60)
    
    # 원본 PPT 복사
    print(f"\n원본 PPT 로드: {ORIGINAL_PPTX}")
    prs = Presentation(ORIGINAL_PPTX)
    original_count = len(prs.slides)
    print(f"기존 슬라이드 수: {original_count}")
    
    # 슬라이드 구조 계획:
    # 1. 표지
    # 2. 분석 대상 소개
    # 3. 연구 목표
    # 4. 데이터 전처리
    # 5. 핵심 인사이트 (토지 구성)
    # 6. 시군구별 비교
    # 7. 대지-공장용지 상관관계
    # 8. 군집별 특성
    # 9. 정책 제언
    # 10-11. 지도 (기존)
    # 12. 참고자료
    
    # 새 슬라이드 추가 (뒤에서부터 추가 후 재정렬)
    print("\n신규 슬라이드 생성 중...")
    
    # 슬라이드 5 이후에 삽입할 새 슬라이드들
    new_slides = []
    
    # 1. 연도별 변화량 (슬라이드 5 다음)
    print("  - 연도별 토지유형 변화량 분석")
    new_slides.append(("yearly_change", create_slide_yearly_change))
    
    # 2. PCA 분류 근거 (슬라이드 7 다음)
    print("  - PCA 분류 근거")
    new_slides.append(("pca_basis", create_slide_pca_basis))
    
    # 3. K-Means 실루엣
    print("  - K-Means 실루엣 분석")
    new_slides.append(("silhouette", create_slide_silhouette))
    
    # 4. 청주 4개구 비교
    print("  - 청주 4개구 비교")
    new_slides.append(("cheongju", create_slide_cheongju))
    
    # 5. 인구 상관관계
    print("  - 인구 상관관계")
    new_slides.append(("pop_correlation", create_slide_pop_correlation))
    
    # 6. 도로 상관관계
    print("  - 도로 상관관계")
    new_slides.append(("road_correlation", create_slide_road_correlation))
    
    # 7. 공장용지 히트맵
    print("  - 공장용지 히트맵")
    new_slides.append(("factory_heatmap", create_slide_factory_heatmap))
    
    # 8. 유형별 추이
    print("  - 유형별 토지 변화 추이")
    new_slides.append(("cluster_trend", create_slide_cluster_trend))
    
    # 슬라이드 생성
    for name, create_func in new_slides:
        create_func(prs)
    
    print(f"\n최종 슬라이드 수: {len(prs.slides)}")
    
    # 저장
    prs.save(OUTPUT_PPTX)
    
    print("\n" + "=" * 60)
    print("PPT 업데이트 완료!")
    print(f"저장 위치: {OUTPUT_PPTX}")
    print("=" * 60)
    
    # 슬라이드 순서 안내
    print("\n※ 슬라이드 순서 조정 필요:")
    print("  현재 새 슬라이드들이 뒤에 추가됨")
    print("  PowerPoint에서 직접 순서 조정 필요:")
    print("  - 슬라이드 13(연도별변화) → 슬라이드 5 다음으로 이동")
    print("  - 슬라이드 14-15(PCA/실루엣) → 슬라이드 7 다음으로 이동")
    print("  - 슬라이드 16(청주4개구) → 슬라이드 6 다음으로 이동")
    print("  - 슬라이드 17-18(인구/도로상관) → 슬라이드 7 다음으로 이동")
    print("  - 슬라이드 19-20(히트맵/추이) → 슬라이드 8 다음으로 이동")


if __name__ == "__main__":
    main()
